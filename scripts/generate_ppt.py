from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "slides" / "rf3-demo-2slides.pptx"
IMG_UI = ROOT / "images" / "rf3-lab1.png"


def _set_font(run, *, size: int, bold: bool = False, color: RGBColor | None = None, name: str = "Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_title(shape, title: str, subtitle: str | None = None):
    tf = shape.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_font(r, size=34, bold=True, color=RGBColor(20, 33, 61))

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = subtitle
        _set_font(r2, size=16, bold=False, color=RGBColor(80, 80, 80))


def add_bullets(shape, bullets: list[str]):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.space_after = Pt(4)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(16)


def build_ppt() -> None:
    prs = Presentation()

    # 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1: Demo 场景
    s1 = prs.slides.add_slide(blank)

    title_box = s1.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.2), Inches(1.3))
    add_title(
        title_box,
        "RosettaFold3 蛋白质实验室 Demo",
        "Azure Container Apps：UI（CPU） + 推理服务（Serverless A100 GPU）",
    )

    left_box = s1.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(6.2), Inches(5.2))
    add_bullets(
        left_box,
        [
            "面向结构生物学 / 蛋白工程 / 药物发现 / 教学演示的在线折叠体验",
            "Web 端提交输入 JSON（支持多链）→ 异步作业（job_id）→ 轮询状态",
            "运行中实时查看 logs tail；完成后下载 predicted.cif/pdb 并网页内渲染（内置 3Dmol.js）",
            "架构：UI FastAPI（反向代理）→ Backend FastAPI（串行执行 rf3 fold / GPU）",
            "在线体验：rf3-demo-ui.bluepebble-ef8ac46c.swedencentral.azurecontainerapps.io",
        ],
    )

    # Screenshot on right
    if IMG_UI.exists():
        s1.shapes.add_picture(str(IMG_UI), Inches(7.3), Inches(2.0), width=Inches(5.7))
    else:
        # Fallback: placeholder box
        ph = s1.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE (avoid enum import)
            Inches(7.3),
            Inches(2.0),
            Inches(5.7),
            Inches(4.6),
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(245, 245, 245)
        ph.line.color.rgb = RGBColor(200, 200, 200)

    # Slide 2: 模型介绍
    s2 = prs.slides.add_slide(blank)

    title2 = s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.2), Inches(1.1))
    add_title(title2, "模型介绍：RosettaFold3（RF3）", "全原子（all-atom）生物大分子结构预测网络")

    bullets2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(7.0), Inches(5.6))
    add_bullets(
        bullets2,
        [
            "输入：氨基酸序列（可选多链/复合物信息）",
            "输出：predicted.cif（mmCIF）/ predicted.pdb（PDB）结构坐标文件（不是图片）",
            "关键建模：隐式手性表示 + 原子级几何条件（atom-level geometric conditioning）",
            "优势任务：手性配体结构/构象预测；固定主链/固定构象条件下 docking",
            "本 Demo 调用：rf3 fold CLI（Hydra 风格 key=value overrides）",
            "部署要点：容器启动下载 latest checkpoint；RF3_CKPT_PATH=/models/rf3.ckpt（建议挂载共享卷做缓存）",
        ],
    )

    side = s2.shapes.add_textbox(Inches(8.2), Inches(1.7), Inches(4.7), Inches(5.6))
    tf = side.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "参考信息"
    _set_font(r, size=18, bold=True, color=RGBColor(20, 33, 61))

    p2 = tf.add_paragraph()
    p2.space_before = Pt(10)
    p2.text = "Preprint：10.1101/2025.08.14.670328"
    p2.font.name = "Microsoft YaHei"
    p2.font.size = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "Checkpoint：rf3_foundry_01_24_latest.ckpt"
    p3.font.name = "Microsoft YaHei"
    p3.font.size = Pt(14)

    p4 = tf.add_paragraph()
    p4.space_before = Pt(8)
    p4.text = "提示：研究/演示用途，结果需实验验证"
    p4.font.name = "Microsoft YaHei"
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(120, 120, 120)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))


if __name__ == "__main__":
    build_ppt()
    print(f"Wrote: {OUT_PATH}")
